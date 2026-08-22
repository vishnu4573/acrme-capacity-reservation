#!/usr/bin/env python3
"""Build the ACRME Executive Presentation (16:9) from the executive design doc.

Output: acrme_executive_presentation.pptx
Content is a faithful executive summary of
`../acrme_executive_design_document.md`.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = os.path.join(os.path.dirname(__file__), "acrme_executive_presentation.pptx")

# ---- Palette --------------------------------------------------------------
AZURE       = RGBColor(0x00, 0x78, 0xD4)   # Azure blue
DEEP        = RGBColor(0x0B, 0x2E, 0x4F)   # deep navy
LIGHT       = RGBColor(0xF3, 0xF7, 0xFB)   # light panel
INK         = RGBColor(0x1F, 0x2A, 0x37)   # body text
MUTED       = RGBColor(0x5B, 0x6B, 0x7B)   # muted text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x1E, 0x8E, 0x3E)   # positive
AMBER       = RGBColor(0xB8, 0x6E, 0x00)   # caution
RED         = RGBColor(0xC5, 0x22, 0x1F)   # blocker
GREY_LINE   = RGBColor(0xD5, 0xDD, 0xE5)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- Helpers --------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def set_run(r, text, size, color, bold=False, italic=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def para(tf, text, size, color, bold=False, italic=False, space_after=6,
         space_before=0, align=PP_ALIGN.LEFT, bullet=False, level=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    r = p.add_run()
    set_run(r, ("•  " + text) if bullet else text, size, color, bold, italic)
    return p


def header(slide, kicker, title, number):
    """Standard content-slide header with a top accent bar."""
    rect(slide, 0, 0, EMU_W, Inches(0.18), AZURE)
    # kicker
    tb, tf = textbox(slide, Inches(0.6), Inches(0.38), Inches(10), Inches(0.35))
    para(tf, kicker.upper(), 12, AZURE, bold=True, first=True, space_after=0)
    # title
    tb, tf = textbox(slide, Inches(0.6), Inches(0.72), Inches(11.6), Inches(0.9))
    para(tf, title, 28, DEEP, bold=True, first=True, space_after=0)
    # rule
    rect(slide, Inches(0.6), Inches(1.55), Inches(12.13), Pt(2), GREY_LINE)
    # page number
    tb, tf = textbox(slide, Inches(12.4), Inches(6.95), Inches(0.8), Inches(0.4))
    para(tf, str(number), 11, MUTED, align=PP_ALIGN.RIGHT, first=True)
    # footer
    tb, tf = textbox(slide, Inches(0.6), Inches(6.95), Inches(9), Inches(0.4))
    para(tf, "ACRME — Executive Design Document  ·  Internal — Restricted",
         9, MUTED, first=True)


def card(slide, x, y, w, h, fill=LIGHT, accent=None):
    if accent:
        rect(slide, x, y, Inches(0.09), h, accent)
    return rect(slide, x, y, w, h, fill)


# ===========================================================================
# SLIDE 1 — Title
# ===========================================================================
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, DEEP)
rect(s, 0, 0, EMU_W, Inches(0.28), AZURE)
rect(s, 0, Inches(4.62), EMU_W, Pt(2), AZURE)
tb, tf = textbox(s, Inches(0.9), Inches(0.7), Inches(8), Inches(0.4))
para(tf, "INTERNAL — RESTRICTED  ·  CONFIDENTIAL", 12, RGBColor(0x9F, 0xC5, 0xE8),
     bold=True, first=True)
tb, tf = textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.0))
para(tf, "Azure Capacity Reservation", 44, WHITE, bold=True, first=True, space_after=0)
para(tf, "Management Engine (ACRME)", 44, WHITE, bold=True, space_after=10)
tb, tf = textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.2))
para(tf, "Executive Design Document — Prepared for Leadership Review", 20,
     RGBColor(0xCF, 0xE2, 0xF3), first=True, space_after=6)
para(tf, "Booking, sharing, and protecting cloud capacity across customers and "
         "regions — with disaster recovery built in.", 14,
     RGBColor(0x9F, 0xC5, 0xE8))
tb, tf = textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
para(tf, "Version 1.0  ·  21 August 2026  ·  Status: Design Complete — Validation In Progress",
     12, RGBColor(0x7F, 0xA9, 0xCC), first=True)

# ===========================================================================
# SLIDE 2 — Executive summary at a glance
# ===========================================================================
s = add_slide()
header(s, "Executive Summary", "The Position in One Slide", 2)
# Left: what it is
tb, tf = textbox(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(3.4))
para(tf, "What ACRME Is", 18, AZURE, bold=True, first=True, space_after=6)
para(tf, "A central control system that books, shares, and protects cloud server "
         "capacity across many customer accounts and locations — a reservation desk "
         "for cloud computing power.", 14, INK, space_after=10)
para(tf, "What It Does", 18, AZURE, bold=True, space_after=6, space_before=4)
for t in [
    "Maintains shared pools of reserved capacity",
    "Chooses the best location for each customer",
    "Protects a dedicated disaster-recovery floor",
    "Requires human approval for high-risk actions",
    "Records every decision for audit",
]:
    para(tf, t, 13, INK, bullet=True, space_after=3)
# Right: status cards
cx, cw = Inches(7.0), Inches(5.7)
cards = [
    ("Design", "Complete & independently reviewed", GREEN),
    ("Fact-check", "7 gaps identified and corrected", GREEN),
    ("Validation", "Planned — not yet executed end-to-end", AMBER),
    ("Production readiness", "≈ 4 / 10 — strong design, unproven execution", RED),
]
cy = Inches(1.8)
for label, val, col in cards:
    card(s, cx, cy, cw, Inches(0.95), LIGHT, accent=col)
    tb, tf = textbox(s, cx + Inches(0.3), cy + Inches(0.1), cw - Inches(0.5),
                     Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, 13, MUTED, bold=True, first=True, space_after=1)
    para(tf, val, 15, DEEP, bold=True)
    cy += Inches(1.12)

# ===========================================================================
# SLIDE 3 — The business problem
# ===========================================================================
s = add_slide()
header(s, "Section 2", "The Business Problem", 3)
tb, tf = textbox(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(1.1))
para(tf, "Cloud capacity is only guaranteed when it is reserved in advance. "
         "Reserved capacity is billed whether or not it is used — so over-reserving "
         "wastes budget and under-reserving creates service risk. Managing that "
         "trade-off manually across dozens of customers and regions does not scale.",
     15, INK, first=True)
tb, tf = textbox(s, Inches(0.6), Inches(3.1), Inches(12.1), Inches(0.4))
para(tf, "What goes wrong without a system", 16, AZURE, bold=True, first=True)
items = [
    "Service disruptions when capacity is unavailable at the moment of need",
    "Failed or severely delayed disaster recoveries — recovery capacity was never secured",
    "Wasted cloud spend from over-reservation nobody is systematically reviewing",
    "Inconsistent practices across teams, making audit and compliance difficult",
    "No single view of what is reserved, used, or protected for recovery",
]
gy = Inches(3.7)
for i, t in enumerate(items):
    card(s, Inches(0.6), gy, Inches(12.13), Inches(0.6), LIGHT, accent=AMBER)
    tb, tf = textbox(s, Inches(0.9), gy, Inches(11.6), Inches(0.6),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 13.5, INK, first=True)
    gy += Inches(0.66)

# ===========================================================================
# SLIDE 4 — The three types of capacity
# ===========================================================================
s = add_slide()
header(s, "Section 3", "The Three Types of Capacity", 4)
tb, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.7))
para(tf, "ACRME keeps three classes of reserved capacity separate in every location "
         "— so everyday demand can never quietly consume what disaster recovery depends on.",
     14, INK, first=True)
cols = [
    ("Production", "Everyday customer workloads that must always run",
     "Always reserved and paid for. Highest day-to-day priority.", AZURE),
    ("Non-production", "Development, testing, and staging",
     "Reserved but flexible. Must not crowd out the DR share.", RGBColor(0x6A,0x4C,0x93)),
    ("Disaster recovery", "Protected floor in a different location",
     "Most protected. Engine-enforced floor. Breaching actions blocked.", GREEN),
]
cw = Inches(3.95)
gap = Inches(0.14)
cx = Inches(0.6)
for name, purpose, treat, col in cols:
    cy = Inches(2.6)
    rect(s, cx, cy, cw, Inches(0.75), col)
    tb, tf = textbox(s, cx, cy, cw, Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, 17, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    rect(s, cx, cy + Inches(0.75), cw, Inches(2.7), LIGHT)
    tb, tf = textbox(s, cx + Inches(0.25), cy + Inches(0.95), cw - Inches(0.5),
                     Inches(2.4))
    para(tf, "Purpose", 11, MUTED, bold=True, first=True, space_after=2)
    para(tf, purpose, 13, INK, space_after=10)
    para(tf, "How it is treated", 11, MUTED, bold=True, space_after=2)
    para(tf, treat, 13, INK)
    cx += cw + gap
tb, tf = textbox(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.8))
para(tf, "Hard rule: a customer's production, non-production, and DR workloads must not "
         "all sit in one geography. Production and DR must be in different locations.",
     13, DEEP, bold=True, first=True)

# ===========================================================================
# SLIDE 5 — What the automation engine does / tiers
# ===========================================================================
s = add_slide()
header(s, "Section 3", "The Automation Engine — Safety Before Automation", 5)
tb, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(6.0), Inches(4.6))
para(tf, "The engine continuously:", 15, AZURE, bold=True, first=True, space_after=6)
for t in [
    "Monitors reservation levels, usage, and headroom",
    "Decides when capacity should be increased",
    "Routes customers to the best location by rules & policy",
    "Prevents any action from destroying the DR floor",
    "Requires human approval for the riskiest operations",
    "Keeps a complete audit trail of every change",
]:
    para(tf, t, 13.5, INK, bullet=True, space_after=5)
# Right: risk tiers
tb, tf = textbox(s, Inches(7.0), Inches(1.75), Inches(5.7), Inches(0.4))
para(tf, "Operations are tiered by risk", 15, AZURE, bold=True, first=True)
tiers = [
    ("Tier 1 — Additive", "Add DR capacity from pre-staged headroom",
     "Allowed in a declared disaster, with fresh checks", GREEN),
    ("Tier 2 — Reallocation", "Reduce non-prod to expand DR",
     "Disabled until validation proves behaviour", AMBER),
    ("Tier 3 — Disassociation", "Change how servers link to reservations",
     "Blocked — needs permissions model + board sign-off", RED),
]
ty = Inches(2.25)
for name, what, treat, col in tiers:
    card(s, Inches(7.0), ty, Inches(5.7), Inches(1.25), LIGHT, accent=col)
    tb, tf = textbox(s, Inches(7.3), ty + Inches(0.12), Inches(5.3), Inches(1.05))
    para(tf, name, 14, col, bold=True, first=True, space_after=2)
    para(tf, what, 12.5, INK, space_after=2)
    para(tf, treat, 12, MUTED, italic=True)
    ty += Inches(1.38)

# ===========================================================================
# SLIDE 6 — Placement decisions (regions)
# ===========================================================================
s = add_slide()
header(s, "Section 6", "How Placement Decisions Are Made", 6)
tb, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.7))
para(tf, "Managed regions are split into two classes. Hard rules filter candidates "
         "first; the engine then scores the survivors on five weighted components.",
     14, INK, first=True)
# Standard regions card
card(s, Inches(0.6), Inches(2.55), Inches(6.0), Inches(2.7), LIGHT, accent=GREEN)
tb, tf = textbox(s, Inches(0.85), Inches(2.7), Inches(5.6), Inches(2.4))
para(tf, "Standard Capacity Regions", 15, GREEN, bold=True, first=True, space_after=6)
for g, r in [("North America", "West US 3, Central US, Canada Central"),
             ("Europe", "Sweden Central, Belgium Central"),
             ("Middle East", "Saudi Arabia, UAE North"),
             ("Asia Pacific", "Japan East, Southeast Asia, Australia East")]:
    p = tf.add_paragraph(); p.space_after = Pt(4)
    r1 = p.add_run(); set_run(r1, g + ":  ", 12.5, DEEP, bold=True)
    r2 = p.add_run(); set_run(r2, r, 12.5, INK)
# Restricted regions card
card(s, Inches(6.73), Inches(2.55), Inches(6.0), Inches(2.7), LIGHT, accent=RED)
tb, tf = textbox(s, Inches(6.98), Inches(2.7), Inches(5.6), Inches(2.4))
para(tf, "Restricted Capacity Regions", 15, RED, bold=True, first=True, space_after=6)
para(tf, "Excluded from automated placement — eligible only via an explicit, "
         "approved exception (production only).", 12.5, INK, space_after=6)
for r in ["East US 2 (North America)", "North Europe, West Europe (Europe)",
          "East Asia, Australia Southeast (Asia Pacific)"]:
    para(tf, r, 12.5, INK, bullet=True, space_after=3)
# Scoring strip
tb, tf = textbox(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(0.4))
para(tf, "Scoring components", 14, AZURE, bold=True, first=True)
tb, tf = textbox(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.9))
para(tf, "Available capacity  ·  Quota headroom  ·  Distribution fairness  ·  "
         "DR coverage health  ·  Zone diversity", 13.5, DEEP, bold=True, first=True,
     space_after=4)
para(tf, "Weights are versioned and recorded with every decision. During the pilot, "
         "placement is recommendation-only — a human commits.", 12.5, MUTED, italic=True)

# ===========================================================================
# SLIDE 7 — Middle East cross-geo
# ===========================================================================
s = add_slide()
header(s, "Section 6", "Middle East — Cross-Geo DR Extension", 7)
tb, tf = textbox(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(1.0))
para(tf, "Middle East has only two Standard Capacity Regions, but a full deployment "
         "needs three (production, non-production, disaster recovery). The engine "
         "extends DR to Belgium Central in Europe — the only currently approved "
         "cross-geography DR path.", 15, INK, first=True)
rows = [
    ("Production", "Higher-scoring of Saudi Arabia or UAE North", AZURE),
    ("Non-production", "The other in-geography region (deterministic)", RGBColor(0x6A,0x4C,0x93)),
    ("Disaster recovery", "Belgium Central (Europe) — cross-geo", GREEN),
]
ry = Inches(3.1)
for env, asg, col in rows:
    card(s, Inches(0.6), ry, Inches(12.13), Inches(0.85), LIGHT, accent=col)
    tb, tf = textbox(s, Inches(0.9), ry, Inches(4.0), Inches(0.85),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, env, 15, col, bold=True, first=True)
    tb, tf = textbox(s, Inches(5.0), ry, Inches(7.5), Inches(0.85),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, asg, 14, INK, first=True)
    ry += Inches(0.98)
tb, tf = textbox(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7))
para(tf, "Belgium Central must still pass capacity and quota checks; if it fails, the "
         "deployment is blocked with an alert — the engine never silently substitutes "
         "another region.", 13, DEEP, bold=True, first=True)

# ===========================================================================
# SLIDE 8 — Disaster recovery
# ===========================================================================
s = add_slide()
header(s, "Section 5", "How Disaster Recovery Works", 8)
tb, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(6.0), Inches(4.5))
para(tf, "Declaration & failover", 16, AZURE, bold=True, first=True, space_after=6)
for t in [
    "Authorised operator requests declaration",
    "Dual approval + state validation required",
    "Engine enters disaster-event mode",
    "Capacity, sharing, zones, quota re-validated",
    "Approved failover deployments begin, observed per workload",
    "Incident-hold mode if state conflicts appear",
]:
    para(tf, t, 13, INK, bullet=True, space_after=4)
tb, tf = textbox(s, Inches(7.0), Inches(1.75), Inches(5.7), Inches(4.5))
para(tf, "The protected capacity floor", 16, AZURE, bold=True, first=True, space_after=6)
para(tf, "A calculated minimum that must stay available for recovery (design "
         "placeholder 30–40%). Enforced by ACRME — not a native cloud sub-reservation.",
     13, INK, space_after=8)
para(tf, "An independent detector recalculates the floor and blocks non-production "
         "expansion if the two calculations disagree (fail-closed).", 13, INK,
     space_after=10)
para(tf, "Failback is gated", 16, AZURE, bold=True, space_after=6)
para(tf, "Controlled return to primary only after explicit approval, readiness "
         "validation, wave-based restore, and health checks. Starting too early is "
         "treated as a critical risk.", 13, INK)

# ===========================================================================
# SLIDE 9 — What's working well
# ===========================================================================
s = add_slide()
header(s, "Section 8", "What Is Working Well", 9)
good = [
    ("Comprehensive, reviewed architecture", "Broad, coherent, independently reviewed; strong requirements coverage"),
    ("Complete risk register", "44 risks logged with likelihood, impact, mitigation, owner, residual"),
    ("Phased approach with human control", "Dangerous operations blocked or approval-gated in Phase 1"),
    ("Strong audit & observability design", "Append-only audit records; dashboards and critical guard alerts"),
    ("Fact-check corrections applied", "7 corrections; high-severity gaps written back into the design"),
    ("Manual survivability", "Operators can still recover if the engine is unavailable"),
]
gx, gy = Inches(0.6), Inches(1.85)
cw, ch = Inches(6.0), Inches(1.5)
for i, (title, body) in enumerate(good):
    col = gx + (cw + Inches(0.13)) * (i % 2)
    row = gy + (ch + Inches(0.13)) * (i // 2)
    card(s, col, row, cw, ch, LIGHT, accent=GREEN)
    tb, tf = textbox(s, col + Inches(0.28), row + Inches(0.15), cw - Inches(0.5),
                     ch - Inches(0.3))
    para(tf, title, 14, DEEP, bold=True, first=True, space_after=3)
    para(tf, body, 12.5, INK)

# ===========================================================================
# SLIDE 10 — What's still being validated (blockers)
# ===========================================================================
s = add_slide()
header(s, "Section 9", "What Is Still Being Validated", 10)
tb, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.5))
para(tf, "None of these should be described as solved until the stated validation is done.",
     14, INK, italic=True, first=True)
blockers = [
    ("Permissions model", "Number-one production blocker — security review & approval required", RED),
    ("End-to-end DR untested", "Failover/failback not yet run live — go-live blocker", RED),
    ("Shared reservation is Preview", "Not covered by GA commitments; can change or be withdrawn", AMBER),
    ("100-customer pool ceiling", "Hard platform limit; pools must be sharded beyond it", AMBER),
    ("Zone alignment across accounts", "Mandatory onboarding gate; needs execution proof per pair", AMBER),
    ("VMSS DR via shared reservation", "Preview limitation; blocked for automated Phase 1 use", AMBER),
]
bx, by = Inches(0.6), Inches(2.4)
cw, ch = Inches(6.0), Inches(1.28)
for i, (title, body, col) in enumerate(blockers):
    cxx = bx + (cw + Inches(0.13)) * (i % 2)
    ryy = by + (ch + Inches(0.13)) * (i // 2)
    card(s, cxx, ryy, cw, ch, LIGHT, accent=col)
    tb, tf = textbox(s, cxx + Inches(0.28), ryy + Inches(0.14), cw - Inches(0.5),
                     ch - Inches(0.28))
    tag = "BLOCKER" if col == RED else "VALIDATING"
    para(tf, tag, 10, col, bold=True, first=True, space_after=2)
    para(tf, title, 14, DEEP, bold=True, space_after=2)
    para(tf, body, 12, INK)

# ===========================================================================
# SLIDE 11 — Risks by category
# ===========================================================================
s = add_slide()
header(s, "Section 10", "The Risks — Plain English", 11)
tb, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.5))
para(tf, "44 logged risks across four dimensions. Highest-impact items and what the "
         "programme is actually doing:", 14, INK, first=True)
risk_cats = [
    ("Technology", [
        "Preview feature changed/withdrawn → feature flags, governance gate",
        "VMSS can't reprovision in outage → block Phase 1, require GA path",
        "Stale inventory drives a bad decision → confirm with direct reads",
    ], AZURE),
    ("Operational", [
        "Permissions model over/under-grants → scoped MI, Tier 3 blocked",
        "DR fails when first exercised → full exercise is a production gate",
        "DR floor bypassed by bug/stale data → independent fail-closed detector",
    ], RGBColor(0x6A,0x4C,0x93)),
    ("Commercial", [
        "Discounts assumed that don't apply → per-customer scope validation",
        "Over-reservation persists in Phase 1 → accepted cost of safety",
        "DR coverage ratio mis-set → 30–40% is placeholder; needs policy",
    ], AMBER),
    ("Customer", [
        "Zone mismatch makes reservation useless → mandatory zone map",
        "100-account ceiling breached → shard pools, monitor continuously",
        "Forced pool removal strands restarts → default-deny + runbook",
    ], GREEN),
]
cx, cy = Inches(0.6), Inches(2.45)
cw, ch = Inches(6.0), Inches(2.05)
for i, (cat, items, col) in enumerate(risk_cats):
    cxx = cx + (cw + Inches(0.13)) * (i % 2)
    ryy = cy + (ch + Inches(0.13)) * (i // 2)
    rect(s, cxx, ryy, cw, Inches(0.5), col)
    tb, tf = textbox(s, cxx + Inches(0.25), ryy, cw - Inches(0.4), Inches(0.5),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, cat + " Risks", 14, WHITE, bold=True, first=True)
    rect(s, cxx, ryy + Inches(0.5), cw, ch - Inches(0.5), LIGHT)
    tb, tf = textbox(s, cxx + Inches(0.25), ryy + Inches(0.6), cw - Inches(0.45),
                     ch - Inches(0.65))
    for t in items:
        para(tf, t, 11.5, INK, bullet=True, space_after=3,
             first=(t == items[0]))

# ===========================================================================
# SLIDE 12 — The three decisions (KEY)
# ===========================================================================
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, DEEP)
rect(s, 0, 0, EMU_W, Inches(0.18), AZURE)
tb, tf = textbox(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35))
para(tf, "SECTION 11 — DECISION", 12, RGBColor(0x9F,0xC5,0xE8), bold=True, first=True)
tb, tf = textbox(s, Inches(0.6), Inches(0.75), Inches(12), Inches(0.7))
para(tf, "Three Decisions Leadership Must Make", 28, WHITE, bold=True, first=True)
decs = [
    ("1", "Accept the Preview feature dependency",
     "Shared reservations are still Microsoft Preview. Recommendation: accept for a "
     "bounded pilot after gates; require a separate production decision later."),
    ("2", "Approve the permissions model approach",
     "The credential model is the #1 production blocker. Recommendation: customer-"
     "consented managed identity with narrow, resource-group-scoped roles — never "
     "subscription-wide admin. Tier 3 stays blocked until security approves."),
    ("3", "Set the disaster recovery floor policy",
     "How much capacity is held for DR is a business decision. 30–40% is a placeholder. "
     "Recommendation: customer-specific recovery analysis for critical workloads."),
]
dy = Inches(1.75)
for num, title, body in decs:
    card(s, Inches(0.6), dy, Inches(12.13), Inches(1.5), RGBColor(0x12,0x3A,0x5E),
         accent=AZURE)
    tb, tf = textbox(s, Inches(0.85), dy, Inches(1.1), Inches(1.5),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, num, 40, AZURE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tb, tf = textbox(s, Inches(2.0), dy + Inches(0.16), Inches(10.5), Inches(1.2))
    para(tf, title, 17, WHITE, bold=True, first=True, space_after=3)
    para(tf, body, 12.5, RGBColor(0xCF,0xE2,0xF3))
    dy += Inches(1.62)
tb, tf = textbox(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.5))
para(tf, "Until each decision is made explicitly, the programme is not cleared for "
         "the corresponding scope.", 12, RGBColor(0x9F,0xC5,0xE8), italic=True, first=True)

# ===========================================================================
# SLIDE 13 — Go-live sequence
# ===========================================================================
s = add_slide()
header(s, "Section 12", "Recommended Go-Live Sequence", 13)
tb, tf = textbox(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.5))
para(tf, "Three stages, each with entry gates. Skipping gates exports risk to customers.",
     14, INK, first=True)
phases = [
    ("Phase 1", "Pilot (Constrained)",
     ["Few customers; manual approval for capacity changes",
      "Placement recommendation-only; DR automation off",
      "Destructive & VMSS emergency actions rejected",
      "Gate: Preview risk accepted, zones validated, rollback proven"], AZURE),
    ("Phase 2", "Controlled Automation",
     ["Automate lower-risk operations on Phase 1 evidence",
      "DR automation remains gated; reallocation only if proven",
      "Tier 3 blocked until permissions + state machine close",
      "Gate: concurrency proven, mode transitions fault-tested"], RGBColor(0x6A,0x4C,0x93)),
    ("Production", "Full Automation (After Gates)",
     ["Only after all gates pass and blockers close",
      "DR tested end-to-end with recovery objectives met",
      "Cost/discount validated per customer vs. billing",
      "Destructive tiers need separate board authorisation"], GREEN),
]
pw = Inches(3.95)
px = Inches(0.6)
for name, sub, items, col in phases:
    py = Inches(2.5)
    rect(s, px, py, pw, Inches(0.95), col)
    tb, tf = textbox(s, px + Inches(0.2), py + Inches(0.12), pw - Inches(0.4),
                     Inches(0.75))
    para(tf, name, 18, WHITE, bold=True, first=True, space_after=0)
    para(tf, sub, 12.5, RGBColor(0xEA,0xF1,0xF8))
    rect(s, px, py + Inches(0.95), pw, Inches(3.1), LIGHT)
    tb, tf = textbox(s, px + Inches(0.25), py + Inches(1.12), pw - Inches(0.5),
                     Inches(2.8))
    for t in items:
        para(tf, t, 12, INK, bullet=True, space_after=6, first=(t == items[0]))
    px += pw + Inches(0.14)
tb, tf = textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
para(tf, "Council position: approve a two-stage programme — constrained pilot now, "
         "controlled production later. Unrestricted autonomous DR is not approved.",
     13, DEEP, bold=True, first=True)

# ===========================================================================
# SLIDE 14 — Immediate actions / close
# ===========================================================================
s = add_slide()
rect(s, 0, 0, EMU_W, EMU_H, DEEP)
rect(s, 0, 0, EMU_W, Inches(0.28), AZURE)
tb, tf = textbox(s, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7))
para(tf, "Immediate Actions for Leadership", 30, WHITE, bold=True, first=True)
actions = [
    "Record a formal Preview-feature decision for the pilot, with a named owner of residual risk",
    "Commission the security review of the customer permissions model and set a decision date",
    "Set interim DR floor and emergency headroom policy values for the pilot cohort",
    "Confirm the pilot customer list and written acceptance of known limitations",
    "Schedule the end-to-end DR exercise as a gated milestone on the production path",
]
ay = Inches(1.7)
for i, t in enumerate(actions, 1):
    card(s, Inches(0.6), ay, Inches(12.13), Inches(0.78), RGBColor(0x12,0x3A,0x5E),
         accent=AZURE)
    tb, tf = textbox(s, Inches(0.85), ay, Inches(0.8), Inches(0.78),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, str(i), 22, AZURE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tb, tf = textbox(s, Inches(1.7), ay, Inches(10.8), Inches(0.78),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, 14, WHITE, first=True)
    ay += Inches(0.9)
rect(s, Inches(0.6), Inches(6.35), Inches(12.13), Pt(2), AZURE)
tb, tf = textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.8))
para(tf, "The honest position: design complete, validation in progress, conditional "
         "pilot supportable — unrestricted autonomous disaster-recovery automation not approved.",
     13.5, RGBColor(0xCF,0xE2,0xF3), italic=True, first=True)

prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
